from io import BytesIO

from docx import Document
from fastapi.testclient import TestClient
from pptx import Presentation
from pypdf import PdfWriter
from uuid import uuid4

from app.auth import get_current_user
from app.ingestion import extract_document
from app.main import app
from app.models import User


def _docx_bytes() -> bytes:
    stream = BytesIO()
    document = Document()
    document.core_properties.title = "Product brief"
    document.add_heading("Problem", level=1)
    document.add_paragraph("Teams spend too much time formatting slides.")
    document.add_heading("Outcome", level=1)
    document.add_paragraph("Generate an editable first draft.")
    document.save(stream)
    return stream.getvalue()


def _pptx_bytes() -> bytes:
    stream = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "MVP\nPilot\nRollout"
    presentation.save(stream)
    return stream.getvalue()


def test_extracts_docx_headings_into_sections() -> None:
    result = extract_document("brief.docx", None, _docx_bytes())
    assert result.title == "Product brief"
    assert [section.title for section in result.sections] == ["Problem", "Outcome"]
    assert "editable first draft" in result.text


def test_extracts_pptx_by_slide() -> None:
    result = extract_document("roadmap.pptx", None, _pptx_bytes())
    assert len(result.sections) == 1
    assert result.sections[0].title == "Roadmap"
    assert "Rollout" in result.text


def test_marks_textless_pdf_for_future_ocr() -> None:
    stream = BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=640, height=480)
    writer.write(stream)

    result = extract_document("scan.pdf", "application/pdf", stream.getvalue())
    assert result.requires_ocr is True
    assert result.warnings


def test_http_contract_accepts_text_and_rejects_unknown_file() -> None:
    user = User(
        id=uuid4(),
        email="member@example.com",
        normalized_email="member@example.com",
        password_hash="not-used",
        is_active=True,
    )
    app.dependency_overrides[get_current_user] = lambda: user
    client = TestClient(app)
    try:
        response = client.post(
            "/v1/ingestion/text",
            json={"kind": "prompt", "title": "Launch", "text": "Create a launch deck"},
        )
        assert response.status_code == 200
        assert response.json()["kind"] == "prompt"

        unsupported = client.post(
            "/v1/ingestion/files",
            files={"file": ("notes.txt", b"not supported", "text/plain")},
        )
        assert unsupported.status_code == 415
    finally:
        app.dependency_overrides.clear()


def test_ingestion_requires_authentication() -> None:
    client = TestClient(app)
    response = client.post(
        "/v1/ingestion/text",
        json={"kind": "prompt", "title": "Private", "text": "Internal only"},
    )
    assert response.status_code == 401
