from io import BytesIO
from pathlib import Path
from zipfile import BadZipFile

from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from pypdf.errors import PdfReadError

from .models import SourceDocument, SourceSection


class UnsupportedDocumentError(ValueError):
    pass


class InvalidDocumentError(ValueError):
    pass


def _clean_text(value: str) -> str:
    lines = [line.rstrip() for line in value.replace("\x00", "").splitlines()]
    return "\n".join(lines).strip()


def _join_sections(sections: list[SourceSection]) -> str:
    return "\n\n".join(section.text for section in sections if section.text).strip()


def _filename_title(filename: str) -> str:
    return Path(filename).stem.strip() or "Untitled document"


def _extract_docx(filename: str, data: bytes) -> SourceDocument:
    try:
        document = Document(BytesIO(data))
    except (BadZipFile, ValueError, KeyError) as error:
        raise InvalidDocumentError("The DOCX file is invalid or corrupted.") from error

    sections: list[SourceSection] = []
    current_title = _filename_title(filename)
    current_lines: list[str] = []

    def flush() -> None:
        nonlocal current_lines
        text = _clean_text("\n".join(current_lines))
        if text:
            sections.append(SourceSection(index=len(sections), title=current_title, text=text))
        current_lines = []

    for paragraph in document.paragraphs:
        text = _clean_text(paragraph.text)
        if not text:
            continue
        if paragraph.style and paragraph.style.name.lower().startswith("heading"):
            flush()
            current_title = text
        else:
            current_lines.append(text)

    for table in document.tables:
        for row in table.rows:
            current_lines.append(" | ".join(_clean_text(cell.text) for cell in row.cells))
    flush()

    title = _clean_text(document.core_properties.title or "") or _filename_title(filename)
    return SourceDocument(kind="docx", title=title, text=_join_sections(sections), sections=sections)


def _shape_text(shape: object) -> str:
    if getattr(shape, "has_text_frame", False):
        return _clean_text(getattr(shape, "text", ""))
    if getattr(shape, "has_table", False):
        table = getattr(shape, "table")
        return "\n".join(
            " | ".join(_clean_text(cell.text) for cell in row.cells) for row in table.rows
        )
    return ""


def _extract_pptx(filename: str, data: bytes) -> SourceDocument:
    try:
        presentation = Presentation(BytesIO(data))
    except (BadZipFile, ValueError, KeyError) as error:
        raise InvalidDocumentError("The PPTX file is invalid or corrupted.") from error

    sections: list[SourceSection] = []
    for slide_index, slide in enumerate(presentation.slides):
        blocks = [text for shape in slide.shapes if (text := _shape_text(shape))]
        title = blocks[0].splitlines()[0][:200] if blocks else f"Slide {slide_index + 1}"
        sections.append(
            SourceSection(index=slide_index, title=title, text=_clean_text("\n".join(blocks)))
        )

    return SourceDocument(
        kind="pptx",
        title=_filename_title(filename),
        text=_join_sections(sections),
        sections=sections,
    )


def _extract_pdf(filename: str, data: bytes) -> SourceDocument:
    try:
        reader = PdfReader(BytesIO(data))
        if reader.is_encrypted and not reader.decrypt(""):
            raise InvalidDocumentError("Password-protected PDFs are not supported.")
        sections = [
            SourceSection(
                index=index,
                title=f"Page {index + 1}",
                text=_clean_text(page.extract_text() or ""),
            )
            for index, page in enumerate(reader.pages)
        ]
    except (PdfReadError, ValueError) as error:
        if isinstance(error, InvalidDocumentError):
            raise
        raise InvalidDocumentError("The PDF file is invalid or corrupted.") from error

    combined = _join_sections(sections)
    requires_ocr = bool(sections) and not combined
    warnings = (
        ["No extractable text was found. Image-only PDF OCR is outside the current MVP."]
        if requires_ocr
        else []
    )
    return SourceDocument(
        kind="pdf",
        title=_filename_title(filename),
        text=combined,
        sections=sections,
        requires_ocr=requires_ocr,
        warnings=warnings,
    )


def extract_document(filename: str, content_type: str | None, data: bytes) -> SourceDocument:
    suffix = Path(filename).suffix.lower()
    if suffix == ".docx":
        return _extract_docx(filename, data)
    if suffix == ".pptx":
        return _extract_pptx(filename, data)
    if suffix == ".pdf":
        return _extract_pdf(filename, data)
    raise UnsupportedDocumentError(
        f"Unsupported document type {suffix or content_type or 'unknown'}; use DOCX, PPTX, or PDF."
    )
